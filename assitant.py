"""
Jarvis Study Assistant (Streamlit app)
Single-file Python app that:
 - Ingests a folder of study materials (pdf, docx, pptx, txt, images, ppt)
 - Extracts text (pdfplumber, python-docx, python-pptx, pytesseract for images)
 - Splits text into chunks
 - Creates embeddings (OpenAI or SentenceTransformers)
 - Stores vectors in Pinecone
 - Answers user queries by retrieving relevant chunks and using an LLM (OpenAI ChatCompletion)

Requirements (pip):
 pip install streamlit langchain openai pinecone-client pdfplumber python-docx python-pptx pillow pytesseract sentence-transformers tiktoken

Environment variables expected:
 - OPENAI_API_KEY (if using OpenAI embeddings/LLM)
 - PINECONE_API_KEY
 - PINECONE_ENV (e.g. us-west1-gcp)
 - PINECONE_INDEX (index name to use/create)
 - EMBEDDING_MODE = "openai" or "sbert" (default: openai)

Notes:
 - You must have Tesseract OCR installed for pytesseract to work (system package).
 - This app is a starting point: tweak chunk size, overlap, and LLM prompt as you like.

Usage:
 streamlit run jarvis_study_assistant.py

"""

import os
import io
import sys
import time
from typing import List, Tuple

import streamlit as st

# Text extraction libraries
import pdfplumber
import docx
from pptx import Presentation
from PIL import Image
import pytesseract

# Vector / embedding / LLM
from sentence_transformers import SentenceTransformer
import numpy as np

# Optional OpenAI
try:
    import openai
except Exception:
    openai = None

# Pinecone
try:
    import pinecone
except Exception:
    pinecone = None

# Small utilities
import hashlib
import json

# -------------------------
# Utility functions
# -------------------------

def file_sha1(path: str) -> str:
    h = hashlib.sha1()
    with open(path, "rb") as f:
        while True:
            b = f.read(8192)
            if not b:
                break
            h.update(b)
    return h.hexdigest()


def read_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def read_pdf(path: str) -> str:
    text_chunks = []
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                txt = page.extract_text() or ""
                text_chunks.append(txt)
    except Exception as e:
        # fallback: try OCR each page image
        try:
            images = pdf2images(path)
            for im in images:
                text_chunks.append(pytesseract.image_to_string(im))
        except Exception:
            print(f"PDF read error: {e}")
    return "\n".join(text_chunks)


def pdf2images(path: str) -> List[Image.Image]:
    # Minimal fallback - requires pdf2image which needs poppler. We avoid mandatory dependency.
    raise RuntimeError("pdf2images not implemented. Install pdf2image and poppler if you need OCR fallback.")


def read_docx(path: str) -> str:
    doc = docx.Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text]
    return "\n".join(paragraphs)


def read_pptx(path: str) -> str:
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text
                if txt:
                    texts.append(txt)
            elif shape.shape_type == 13:  # picture
                try:
                    img = shape.image
                except Exception:
                    img = None
        # Additionally try to OCR entire slide if needed (not implemented here)
    return "\n".join(texts)


def read_image(path: str) -> str:
    im = Image.open(path)
    txt = pytesseract.image_to_string(im)
    return txt


def extract_text_from_file(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext in [".txt"]:
        return read_txt(path)
    elif ext in [".pdf"]:
        return read_pdf(path)
    elif ext in [".docx"]:
        return read_docx(path)
    elif ext in [".pptx", ".ppt"]:
        return read_pptx(path)
    elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff"]:
        return read_image(path)
    else:
        # try reading as text
        try:
            return read_txt(path)
        except Exception:
            return ""


# Simple text splitter

def split_text(text: str, chunk_size: int = 800, overlap: int = 200) -> List[str]:
    if not text:
        return []
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = words[i : i + chunk_size]
        chunks.append(" ".join(chunk))
        i += chunk_size - overlap
    return chunks


# Embedding management
class Embedder:
    def __init__(self, mode: str = "openai"):
        self.mode = mode
        if mode == "sbert":
            self.model = SentenceTransformer("all-MiniLM-L6-v2")
        elif mode == "openai":
            if openai is None:
                raise RuntimeError("openai package not available")
            # Key should be in env var OPENAI_API_KEY
            openai.api_key = os.getenv("OPENAI_API_KEY")
        else:
            raise ValueError("Unknown embedding mode")

    def embed(self, texts: List[str]) -> List[List[float]]:
        if self.mode == "sbert":
            embs = self.model.encode(texts, show_progress_bar=False, convert_to_numpy=True)
            return [e.tolist() for e in embs]
        else:
            # OpenAI embeddings
            # using "text-embedding-3-small" or "text-embedding-3-large"
            embs = []
            BATCH = 32
            for i in range(0, len(texts), BATCH):
                batch = texts[i : i + BATCH]
                resp = openai.Embedding.create(input=batch, model="text-embedding-3-small")
                for r in resp["data"]:
                    embs.append(r["embedding"])
            return embs


# Pinecone helper
class PineconeStore:
    def __init__(self, index_name: str, dimension: int, namespace: str = None):
        if pinecone is None:
            raise RuntimeError("pinecone package not installed")
        api_key = os.getenv("PINECONE_API_KEY")
        env = os.getenv("PINECONE_ENV")
        if not api_key or not env:
            raise RuntimeError("PINECONE_API_KEY and PINECONE_ENV must be set in environment variables")
        pinecone.init(api_key=api_key, environment=env)
        self.index_name = index_name
        self.dimension = dimension
        self.namespace = namespace
        if index_name not in pinecone.list_indexes():
            pinecone.create_index(index_name, dimension=dimension)
            # Wait for index creation
            time.sleep(2)
        self.index = pinecone.Index(index_name)

    def upsert(self, vectors: List[Tuple[str, List[float], dict]]):
        # vectors: list of (id, vector, metadata)
        # Pinecone upsert accepts list of tuples
        self.index.upsert(vectors=vectors, namespace=self.namespace)

    def query(self, vector: List[float], top_k: int = 5):
        res = self.index.query(vector=vector, top_k=top_k, include_metadata=True, namespace=self.namespace)
        return res


# LLM responder using OpenAI ChatCompletion
class Responder:
    def __init__(self):
        if openai is None:
            raise RuntimeError("openai package not available")
        openai.api_key = os.getenv("OPENAI_API_KEY")

    def answer(self, question: str, contexts: List[dict]) -> str:
        # contexts: list of metadata dicts containing 'text' and 'source'
        system_prompt = "You are Jarvis, a helpful assistant. Use the provided context excerpts to answer the user's question. If you don't know, say you don't know. Keep answers concise and cite source filenames when relevant."
        # Build a context string
        context_texts = []
        for i, c in enumerate(contexts):
            src = c.get("source", "unknown")
            txt = c.get("text", "")
            excerpt = txt[:800].strip().replace("\n", " ")
            context_texts.append(f"[{i+1}] Source: {src}\n{excerpt}")
        prompt = system_prompt + "\n\nCONTEXT:\n" + "\n---\n".join(context_texts) + "\n\nUSER QUESTION:\n" + question

        resp = openai.ChatCompletion.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt},
            ],
            max_tokens=512,
            temperature=0.0,
        )
        text = resp["choices"][0]["message"]["content"].strip()
        return text


# -------------------------
# Streamlit UI + Main flows
# -------------------------

st.set_page_config(page_title="Jarvis - Study Assistant", layout="wide")
st.title("Jarvis — Personal Study Assistant")
st.write("Ingest your study folder, store vectors in Pinecone, and ask questions grounded on your files.")

# Sidebar: configuration
st.sidebar.header("Configuration")
upload_folder = st.sidebar.text_input("Folder path with study files", value="./study_materials")
index_name = st.sidebar.text_input("Pinecone index name", value=os.getenv("PINECONE_INDEX", "jarvis-index"))
embedding_mode = st.sidebar.selectbox("Embedding mode", options=["openai", "sbert"], index=0 if os.getenv("EMBEDDING_MODE","openai")=="openai" else 1)
chunk_size = st.sidebar.number_input("Chunk size (words)", value=800, min_value=100, max_value=4000, step=100)
overlap = st.sidebar.number_input("Chunk overlap (words)", value=200, min_value=0, max_value=1000, step=50)
namespace = st.sidebar.text_input("Pinecone namespace (optional)", value="jarvis")

st.sidebar.markdown("---")
if st.sidebar.button("Ingest folder to Pinecone"):
    with st.spinner("Ingesting files — extracting text, embedding, and uploading to Pinecone..."):
        all_files = []
        for root, dirs, files in os.walk(upload_folder):
            for f in files:
                if f.startswith("."):
                    continue
                path = os.path.join(root, f)
                all_files.append(path)

        st.write(f"Found {len(all_files)} files.")
        embedder = Embedder(mode=embedding_mode)
        # Determine embedding dim
        sample_emb = None
        try:
            sample_emb = embedder.embed(["hello world"])[0]
        except Exception as e:
            st.error(f"Embedding test failed: {e}")
            st.stop()
        dim = len(sample_emb)
        # Init Pinecone store
        store = PineconeStore(index_name=index_name, dimension=dim, namespace=namespace)

        upsert_batch = []
        BATCH = 100
        total_chunks = 0
        for path in all_files:
            try:
                text = extract_text_from_file(path)
            except Exception as e:
                st.warning(f"Failed extract {path}: {e}")
                continue
            if not text:
                continue
            chunks = split_text(text, chunk_size=chunk_size, overlap=overlap)
            total_chunks += len(chunks)
            embeddings = embedder.embed(chunks)
            for i, (chunk, emb) in enumerate(zip(chunks, embeddings)):
                uid = hashlib.sha1((path + str(i)).encode()).hexdigest()
                metadata = {"source": os.path.relpath(path, upload_folder), "text": chunk[:10000]}
                upsert_batch.append((uid, emb, metadata))
                # upsert in batches
                if len(upsert_batch) >= BATCH:
                    store.upsert(upsert_batch)
                    upsert_batch = []
        if upsert_batch:
            store.upsert(upsert_batch)
        st.success(f"Ingested {total_chunks} chunks into Pinecone index '{index_name}' (namespace: {namespace}).")

st.sidebar.markdown("---")

# Query UI
st.header("Ask questions from your ingested files")
query = st.text_area("Enter question", height=120)
top_k = st.number_input("Top K retrievals", value=4, min_value=1, max_value=10)

col1, col2 = st.columns([1, 3])
with col1:
    if st.button("Search & Answer"):
        if not query.strip():
            st.warning("Please type a question first")
        else:
            with st.spinner("Embedding question and querying Pinecone..."):
                embedder = Embedder(mode=embedding_mode)
                q_emb = embedder.embed([query])[0]
                store = PineconeStore(index_name=index_name, dimension=len(q_emb), namespace=namespace)
                res = store.query(q_emb, top_k=top_k)
                # Parse hits
                hits = res.get("matches", [])
                contexts = []
                for h in hits:
                    meta = h.get("metadata", {})
                    contexts.append({"score": h.get("score"), "source": meta.get("source","unknown"), "text": meta.get("text","")})

            # Use LLM to answer
            with st.spinner("Calling LLM to generate answer..."):
                responder = Responder()
                answer = responder.answer(query, contexts)

            st.markdown("**Answer:**")
            st.write(answer)
            st.markdown("**Retrieved sources (top results):**")
            for i, c in enumerate(contexts):
                st.write(f"{i+1}. {c.get('source')} — score: {c.get('score'):.4f}")
                st.text(c.get('text')[:800])

with col2:
    st.markdown("## Quick instructions")
    st.markdown(
        "- Place your files (pdf/docx/pptx/txt/images) in the folder you set in the sidebar.\n- Click `Ingest folder to Pinecone`.\n- Wait for ingestion to finish.\n- Ask questions in the box and press `Search & Answer`.\n- Configure embedding mode (openai or sbert).\n"
    )
    st.markdown("### Environment variables needed")
    st.code("""
export OPENAI_API_KEY=sk-...
export PINECONE_API_KEY=...
export PINECONE_ENV=us-west1-gcp
export PINECONE_INDEX=jarvis-index
export EMBEDDING_MODE=openai
""")

st.markdown("---")
st.markdown("Built as a starting point. You can extend: more file types, better OCR, metadata extraction, caching embeddings locally, alternative LLMs, fallback strategies.")

# Save settings button
if st.button("Save config to local file"):
    cfg = {"upload_folder": upload_folder, "index_name": index_name, "embedding_mode": embedding_mode, "chunk_size": chunk_size, "overlap": overlap, "namespace": namespace}
    with open("jarvis_config.json", "w") as f:
        json.dump(cfg, f, indent=2)
    st.success("Config saved to jarvis_config.json")